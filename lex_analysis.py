import os
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime

input_path = "data/analysis_input.csv"
output_dir = "output"
os.makedirs(output_dir, exist_ok=True)

try:
    df = pd.read_csv(input_path)
    print("[LEX] Daten erfolgreich geladen.")
except FileNotFoundError:
    print(f"[LEX] Fehler: Datei '{input_path}' nicht gefunden.")
    exit(1)

print("[LEX] Statistische Auswertung:")
print(df.describe())

plt.figure(figsize=(10, 6))
df.hist()
plt.tight_layout()
hist_path = os.path.join(output_dir, "analysis_histogram.png")
plt.savefig(hist_path)
plt.close()
print(f"[LEX] Diagramm gespeichert: {hist_path}")

bericht_path = os.path.join(output_dir, "bericht.txt")
with open(bericht_path, "w", encoding="utf-8") as f:
    f.write("QuantumShield  Analysebericht\\n")
    f.write("================================\\n")
    f.write(f"Datum: {datetime.now().strftime('%d.%m.%Y %H:%M:%S')}\\n")
    f.write(f"Datenquelle: {input_path}\\n")
    f.write("Analyseergebnisse (describe):\\n\\n")
    f.write(df.describe().to_string())

print(f"[LEX] Bericht gespeichert: {bericht_path}")

ppt_path = os.path.join(output_dir, "QuantumShield_Praesentation.pptx")
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Analysebericht  QuantumShield"
slide.placeholders[1].text = f"Erstellt am {datetime.now().strftime('%d.%m.%Y')}"

slide2 = prs.slides.add_slide(prs.slide_layouts[5])
slide2.shapes.title.text = "Histogramm"
slide2.shapes.add_picture(hist_path, Inches(1), Inches(1.5), width=Inches(6))

prs.save(ppt_path)
print(f"[LEX] Präsentation gespeichert: {ppt_path}")
